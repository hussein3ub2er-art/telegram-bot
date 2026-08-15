"""
وحدة التحويل المحلي بين الملفات (بدون إنترنت)
------------------------------------------------
كل التحويلات هنا "نص فقط" — تستخرج النص وتعيد بناءه بالصيغة الجديدة،
بدون صور أو تنسيقات معقدة أو تصميم أصلي. هذا لأن التحويل الكامل
(بنفس الشكل والتصميم) يحتاج LibreOffice أو PowerPoint، وهذي البرامج
غير متوفرة على الجوال.

المكتبات المستخدمة كلها خفيفة ونقية بايثون (ما تحتاج ترجمة/compile):
pypdf, python-docx, python-pptx, reportlab
"""

from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm

MAX_CHARS_PER_LINE = 95
LINE_HEIGHT = 14


def _wrap(text: str) -> list[str]:
    text = text.strip()
    if not text:
        return [""]
    return [
        text[i : i + MAX_CHARS_PER_LINE]
        for i in range(0, len(text), MAX_CHARS_PER_LINE)
    ]


def _text_to_pdf(pages_text: list[str], pdf_path: str) -> None:
    """كتابة قائمة نصوص (نص لكل صفحة/شريحة) في ملف PDF بسيط."""
    c = canvas.Canvas(pdf_path, pagesize=A4)
    width, height = A4
    margin = 2 * cm

    for page_text in pages_text:
        y = height - margin
        c.setFont("Helvetica", 11)
        for raw_line in page_text.split("\n"):
            for line in _wrap(raw_line):
                if y < margin:
                    c.showPage()
                    c.setFont("Helvetica", 11)
                    y = height - margin
                c.drawString(margin, y, line)
                y -= LINE_HEIGHT
        c.showPage()

    c.save()


def pdf_to_docx(pdf_path: str, docx_path: str) -> None:
    """استخراج نص PDF وكتابته في ملف Word (نص فقط، بدون صور أو تنسيق)."""
    reader = PdfReader(pdf_path)
    document = Document()

    for page_number, page in enumerate(reader.pages, start=1):
        if page_number > 1:
            document.add_page_break()
        document.add_heading(f"صفحة {page_number}", level=2)
        text = page.extract_text() or ""
        for paragraph in text.split("\n"):
            document.add_paragraph(paragraph)

    document.save(docx_path)


def docx_to_pdf(docx_path: str, pdf_path: str) -> None:
    """تحويل Word إلى PDF (نص فقط، بدون صور أو تنسيقات معقدة)."""
    document = Document(docx_path)
    full_text = "\n".join(para.text for para in document.paragraphs)
    _text_to_pdf([full_text], pdf_path)


def pdf_to_pptx(pdf_path: str, pptx_path: str) -> None:
    """استخراج نص كل صفحة PDF ووضعه في شريحة منفصلة (نص فقط)."""
    reader = PdfReader(pdf_path)
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for page_number, page in enumerate(reader.pages, start=1):
        slide = prs.slides.add_slide(blank_layout)
        text_box = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        title_para = text_frame.paragraphs[0]
        title_para.text = f"صفحة {page_number}"

        body_text = page.extract_text() or ""
        for line in body_text.split("\n"):
            if not line.strip():
                continue
            p = text_frame.add_paragraph()
            p.text = line

    prs.save(pptx_path)


def pptx_to_pdf(pptx_path: str, pdf_path: str) -> None:
    """تحويل نص شرائح PowerPoint إلى PDF (نص فقط، بدون صور أو تصميم)."""
    prs = Presentation(pptx_path)
    pages_text = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        lines = [f"شريحة {slide_number}", ""]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    lines.append(text)
        pages_text.append("\n".join(lines))

    _text_to_pdf(pages_text, pdf_path)
